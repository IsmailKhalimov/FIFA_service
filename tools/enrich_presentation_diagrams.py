"""
Добавляет в pptx два слайда со схемами связей (вместо пяти): стек + потоки.

Вставка: после слайда «Сущности и роли», перед «Интерфейс: три экрана».

Запуск из корня:
  python tools/enrich_presentation_diagrams.py
  python tools/enrich_presentation_diagrams.py report/presentation_with_diagrams.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

import matplotlib.image as mpimg
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
REPORT = ROOT / "report"
DEFAULT_PPTX = REPORT / "presentation.pptx"

BG = RGBColor(14, 17, 23)
TITLE = RGBColor(250, 250, 250)

MARKER_PREFIX = "Схема: "


def move_slide(presentation: Presentation, old_index: int, new_index: int) -> None:
    sld_id_lst = presentation.slides._sldIdLst  # noqa: SLF001
    ids = list(sld_id_lst)
    el = ids[old_index]
    sld_id_lst.remove(el)
    sld_id_lst.insert(new_index, el)


def _style_axes_dark(ax) -> None:
    ax.set_facecolor("#161a22")
    for spine in ax.spines.values():
        spine.set_color("#3a3f4b")


def _setup_matplotlib_ru() -> None:
    rcParams["font.family"] = "Segoe UI"
    rcParams["figure.facecolor"] = "#0e1117"
    rcParams["text.color"] = "#f0f2f5"
    rcParams["axes.edgecolor"] = "#3a3f4b"


def _box(ax, xy, w, h, text, fc="#2d3340", ec="#ff4b4b", fontsize=11):
    box = mpatches.FancyBboxPatch(
        xy, w, h, boxstyle="round,pad=0.02,rounding_size=0.08",
        linewidth=2, edgecolor=ec, facecolor=fc,
    )
    ax.add_patch(box)
    ax.text(xy[0] + w / 2, xy[1] + h / 2, text, ha="center", va="center", fontsize=fontsize, color="#f0f2f5", wrap=True)


def _arrow(ax, x1, y1, x2, y2):
    ax.annotate(
        "", xy=(x2, y2), xytext=(x1, y1),
        arrowprops=dict(arrowstyle="->", color="#8b92a8", lw=2),
    )


def diagram_context(path: Path) -> None:
    _setup_matplotlib_ru()
    fig, ax = plt.subplots(figsize=(12, 6.2), dpi=140)
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 6)
    ax.axis("off")
    _style_axes_dark(ax)
    ax.set_title("Клиент → приложение → БД", color="#f0f2f5", fontsize=16, pad=12)
    boxes = [
        (0.3, 2.4, 1.6, 1.1, "Браузер"),
        (2.5, 2.4, 1.8, 1.1, "Streamlit"),
        (4.8, 2.4, 1.7, 1.1, "Python"),
        (7.0, 2.4, 1.7, 1.1, "psycopg2"),
        (9.2, 2.4, 2.4, 1.1, "PostgreSQL"),
    ]
    for b in boxes:
        _box(ax, (b[0], b[1]), b[2], b[3], b[4], fontsize=12)
    for i in range(len(boxes) - 1):
        x1 = boxes[i][0] + boxes[i][2]
        y = boxes[i][1] + boxes[i][3] / 2
        x2 = boxes[i + 1][0]
        _arrow(ax, x1, y, x2, y)
    _box(ax, (3.8, 0.35), 4.6, 0.95, "Фото: отдельный скрипт → URL в БД", fc="#1e2430", ec="#6b7280", fontsize=11)
    _arrow(ax, 6.1, 1.3, 6.1, 2.35)
    fig.text(0.5, 0.04, "Рабочие запросы идут в PostgreSQL", ha="center", fontsize=11, color="#9ca3af")
    plt.tight_layout(rect=[0, 0.06, 1, 1])
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def diagram_modules(path: Path) -> None:
    _setup_matplotlib_ru()
    fig, ax = plt.subplots(figsize=(12, 6.4), dpi=140)
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 7)
    ax.axis("off")
    _style_axes_dark(ax)
    ax.set_title("Файлы проекта", color="#f0f2f5", fontsize=16, pad=12)
    _box(ax, (4.9, 5.5), 2.2, 0.85, "main.py", fontsize=12)
    _box(ax, (4.7, 4.0), 2.6, 0.9, "pages.py", fontsize=12)
    _box(ax, (1.0, 2.2), 2.5, 0.9, "auth.py", fontsize=12)
    _box(ax, (4.7, 2.2), 2.6, 0.9, "db_connection.py", fontsize=12)
    _box(ax, (8.4, 2.2), 2.5, 0.9, "agg_func.py", fontsize=12)
    _box(ax, (8.4, 4.0), 2.5, 0.9, "tm_images.py", fontsize=12)
    _box(ax, (4.7, 0.45), 2.6, 0.95, "config.py", fontsize=12)
    _arrow(ax, 6.0, 5.5, 6.0, 4.9)
    _arrow(ax, 5.5, 4.0, 3.5, 3.1)
    _arrow(ax, 6.0, 4.0, 6.0, 3.1)
    _arrow(ax, 6.5, 4.0, 9.0, 3.1)
    _arrow(ax, 6.5, 4.0, 9.0, 4.5)
    _arrow(ax, 6.0, 2.2, 6.0, 1.4)
    fig.text(0.5, 0.02, "Стрелки: типичный вызов при действии пользователя", ha="center", fontsize=11, color="#9ca3af")
    plt.tight_layout(rect=[0, 0.05, 1, 1])
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def diagram_filter_flow(path: Path) -> None:
    _setup_matplotlib_ru()
    fig, ax = plt.subplots(figsize=(12, 5.8), dpi=140)
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 5.5)
    ax.axis("off")
    _style_axes_dark(ax)
    ax.set_title("Фильтрация клубов", color="#f0f2f5", fontsize=16, pad=12)
    row_y = 2.0
    _box(ax, (0.4, row_y), 2.0, 1.0, "Фильтры\nв UI", fontsize=12)
    _box(ax, (2.8, row_y), 2.2, 1.0, "Условия\nWHERE", fontsize=12)
    _box(ax, (5.4, row_y), 2.4, 1.0, "get_filtered_clubs()", fontsize=11)
    _box(ax, (8.2, row_y), 3.2, 1.0, "Таблица на экране", fontsize=12)
    _arrow(ax, 2.4, row_y + 0.5, 2.8, row_y + 0.5)
    _arrow(ax, 5.0, row_y + 0.5, 5.4, row_y + 0.5)
    _arrow(ax, 7.8, row_y + 0.5, 8.2, row_y + 0.5)
    fig.text(0.5, 0.08, "Форма → SQL → результат", ha="center", fontsize=11, color="#9ca3af")
    plt.tight_layout(rect=[0, 0.06, 1, 1])
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def diagram_recommend_flow(path: Path) -> None:
    _setup_matplotlib_ru()
    fig, ax = plt.subplots(figsize=(12, 6.0), dpi=140)
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 6)
    ax.axis("off")
    _style_axes_dark(ax)
    ax.set_title("Рекомендации (вариант B)", color="#f0f2f5", fontsize=16, pad=12)
    y = 2.2
    _box(ax, (0.35, y), 2.3, 1.05, "Форма:\nцель и веса", fontsize=11)
    _box(ax, (3.0, y), 2.6, 1.05, "search_ideal_profile_players()", fontsize=10)
    _box(ax, (6.0, y), 2.5, 1.05, "ORDER BY\nfit_dist", fontsize=12)
    _box(ax, (9.0, y), 2.7, 1.05, "Список\nигроков", fontsize=12)
    _arrow(ax, 2.65, y + 0.52, 3.0, y + 0.52)
    _arrow(ax, 5.6, y + 0.52, 6.0, y + 0.52)
    _arrow(ax, 8.5, y + 0.52, 9.0, y + 0.52)
    _box(ax, (2.5, 0.35), 7.0, 1.05, "fit_dist — взвешенные отклонения от цели (как в SQL)", fc="#1e2430", ec="#6b7280", fontsize=11)
    _arrow(ax, 6.0, 1.4, 6.0, 2.15)
    fig.text(0.5, 0.04, "Те же кандидаты, что у скаута; другая сортировка", ha="center", fontsize=11, color="#9ca3af")
    plt.tight_layout(rect=[0, 0.06, 1, 1])
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def combine_vertical(top: Path, bottom: Path, out: Path) -> None:
    _setup_matplotlib_ru()
    fig, axes = plt.subplots(2, 1, figsize=(12, 11.2), dpi=140)
    fig.patch.set_facecolor("#0e1117")
    for ax, p in zip(axes, (top, bottom), strict=True):
        im = mpimg.imread(str(p))
        ax.imshow(im)
        ax.axis("off")
    plt.subplots_adjust(hspace=0.08, left=0.02, right=0.98, top=0.99, bottom=0.01)
    fig.savefig(out, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def _set_slide_dark_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.82))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TITLE
    p.alignment = PP_ALIGN.LEFT


def _speaker(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text.strip()


def _add_diagram_slide(prs: Presentation, title: str, img: Path, notes: str) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_slide_dark_bg(slide)
    _add_title(slide, title)
    slide.shapes.add_picture(str(img), Inches(0.4), Inches(0.95), height=Inches(5.55))
    _speaker(slide, notes)


def _already_enriched(prs: Presentation) -> bool:
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t.startswith(MARKER_PREFIX) or t.startswith("Диаграмма: контекст"):
                    return True
    return False


def main() -> None:
    pptx_arg = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_PPTX
    pptx_path = pptx_arg if pptx_arg.is_absolute() else ROOT / pptx_arg
    if not pptx_path.exists():
        raise SystemExit(f"Не найден файл: {pptx_path}")

    REPORT.mkdir(parents=True, exist_ok=True)
    p_ctx = REPORT / "pres_diag_01_context.png"
    p_mod = REPORT / "pres_diag_02_modules.png"
    p_flt = REPORT / "pres_diag_03_filter_flow.png"
    p_rec = REPORT / "pres_diag_04_recommend_flow.png"
    diagram_context(p_ctx)
    diagram_modules(p_mod)
    diagram_filter_flow(p_flt)
    diagram_recommend_flow(p_rec)

    bundle_a = REPORT / "pres_diag_bundle_system.png"
    bundle_b = REPORT / "pres_diag_bundle_flows.png"
    combine_vertical(p_ctx, p_mod, bundle_a)
    combine_vertical(p_flt, p_rec, bundle_b)

    prs = Presentation(str(pptx_path))
    if _already_enriched(prs):
        print("Схемы уже есть в презентации — пропуск.")
        return

    n_before = len(prs.slides)
    # После «Сущности и роли» (5-й слайд, индекс 4) идёт UI — вставляем с индекса 5
    insert_at = 5

    pairs = [
        (
            MARKER_PREFIX + "клиент и модули (две схемы)",
            bundle_a,
            "Верх: цепочка до БД. Низ: какие файлы за что отвечают. 30–40 с.",
        ),
        (
            MARKER_PREFIX + "фильтр клубов и рекомендации",
            bundle_b,
            "Верх: фильтр → SQL → таблица. Низ: форма → тот же SQL с fit_dist → список.",
        ),
    ]

    for title, img, notes in pairs:
        _add_diagram_slide(prs, title, Path(img), notes)

    for k in range(2):
        move_slide(prs, n_before, insert_at + k)

    out = pptx_path
    try:
        prs.save(str(out))
    except PermissionError:
        alt = REPORT / "presentation_with_diagrams.pptx"
        prs.save(str(alt))
        print(f"{out.name} занят. Сохранено: {alt}")
    else:
        print(f"Добавлено 2 слайда со схемами. Всего слайдов: {len(prs.slides)}. Файл: {out}")


if __name__ == "__main__":
    main()
