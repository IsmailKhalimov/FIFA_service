"""
Сборка презентации для защиты: report/presentation.pptx
Расчёт ~7 минут: меньше слайдов, крупнее шрифт, короче формулировки.

Запуск из корня репозитория:
  pip install python-pptx matplotlib
  python tools/build_presentation.py
"""
from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib import rcParams
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
REPORT = ROOT / "report"
ASSETS = Path(r"C:\Users\lirik\.cursor\projects\d-gits-DataBase-MAI-main\assets")

BG = RGBColor(14, 17, 23)
TITLE = RGBColor(250, 250, 250)
BODY = RGBColor(220, 224, 230)
ACCENT = RGBColor(255, 75, 75)

TITLE_PT = 40
BODY_PT = 26
IMG_TITLE_PT = 36


def _set_slide_dark_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _add_title(
    slide,
    text: str,
    *,
    top: float = 0.32,
    height: float = 0.82,
    size_pt: int = TITLE_PT,
) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.2), Inches(height))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(size_pt)
    p.font.bold = True
    p.font.color.rgb = TITLE
    p.alignment = PP_ALIGN.LEFT


def _add_body(slide, lines: list[str], top: float = 1.12, width: float = 12.2, size_pt: int = BODY_PT) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(width), Inches(5.95))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Segoe UI"
        p.font.size = Pt(size_pt)
        p.font.color.rgb = BODY
        p.space_after = Pt(8)
        p.alignment = PP_ALIGN.LEFT


def _speaker(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = text.strip()


def _fit_dist(
    age: float,
    salary: float,
    goals: float | None,
    assists: float | None,
    *,
    t_age: float = 26.0,
    t_sal: float = 8_000_000.0,
    t_goals: float = 12.0,
    t_assists: float = 8.0,
    w_age: float = 1.0,
    w_sal: float = 1.0,
    w_goals: float = 1.2,
    w_assists: float = 1.0,
) -> float:
    g = 0.0 if goals is None else float(goals)
    a = 0.0 if assists is None else float(assists)
    return (
        w_age * abs(age - t_age) / 40.0
        + w_sal * abs(salary - t_sal) / 50_000_000.0
        + w_goals * abs(g - t_goals) / 30.0
        + w_assists * abs(a - t_assists) / 25.0
    )


def _build_quality_chart(path: Path) -> None:
    rcParams["font.family"] = "Segoe UI"
    rcParams["axes.facecolor"] = "#161a22"
    rcParams["figure.facecolor"] = "#0e1117"
    rcParams["text.color"] = "#f0f2f5"
    rcParams["axes.labelcolor"] = "#f0f2f5"
    rcParams["xtick.color"] = "#f0f2f5"
    rcParams["ytick.color"] = "#f0f2f5"
    rcParams["axes.edgecolor"] = "#3a3f4b"

    players = [
        ("Идеальный кандидат", 26, 8_000_000, 12, 8),
        ("Почти совпадение", 25, 7_800_000, 13, 7),
        ("Сильнее в голах", 27, 8_200_000, 18, 6),
        ("Дешевле и моложе", 24, 5_500_000, 10, 9),
        ("Смещение по зарплате", 26, 12_000_000, 11, 8),
        ("Мало передач", 26, 8_000_000, 12, 4),
        ("Старше профиля", 31, 8_000_000, 11, 9),
        ("Высокая зарплата", 26, 15_000_000, 10, 8),
        ("Низкая результативность", 26, 8_000_000, 4, 5),
        ("Далёкий профиль", 22, 3_000_000, 5, 3),
        ("Очень далёкий профиль", 34, 18_000_000, 3, 2),
    ]
    scored = [(name, _fit_dist(age, sal, g, a)) for name, age, sal, g, a in players]
    scored.sort(key=lambda x: x[1])
    names = [s[0] for s in scored]
    vals = [s[1] for s in scored]

    fig, ax = plt.subplots(figsize=(11, 6.4), dpi=140)
    y_pos = range(len(names))
    bars = ax.barh(list(y_pos), vals, color="#ff4b4b", height=0.65, alpha=0.9)
    ax.set_yticks(list(y_pos))
    ax.set_yticklabels(names, fontsize=13)
    ax.invert_yaxis()
    ax.set_xlabel("Расстояние до профиля (меньше — ближе)", fontsize=14)
    ax.set_title("Ранжирование по той же формуле, что в SQL", fontsize=15, pad=14)
    for bar, v in zip(bars, vals):
        ax.text(
            v + 0.02,
            bar.get_y() + bar.get_height() / 2,
            f"{v:.3f}",
            va="center",
            ha="left",
            color="#f0f2f5",
            fontsize=11,
        )
    ax.grid(axis="x", color="#3a3f4b", linestyle="--", alpha=0.6)
    fig.text(
        0.5,
        0.02,
        "Цель: возраст 26, зарплата 8 млн $, голы 12, передачи 8; веса 1 / 1 / 1,2 / 1",
        ha="center",
        fontsize=12,
        color="#b8bcc8",
    )
    plt.tight_layout(rect=[0, 0.06, 1, 1])
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def _combine_three_screenshots(slide, paths: list[Path]) -> None:
    """Три скрина в один ряд, фиксированная ширина (пропорции сохраняются)."""
    top = Inches(1.08)
    w = Inches(4.0)
    gap = Inches(0.14)
    left0 = Inches(0.38)
    step = int(w) + int(gap)
    base = int(left0)
    for i, p in enumerate(paths):
        slide.shapes.add_picture(str(p), base + i * step, top, width=w)


def main() -> None:
    REPORT.mkdir(parents=True, exist_ok=True)
    shots = [
        ASSETS
        / "c__Users_lirik_AppData_Roaming_Cursor_User_workspaceStorage_001952add7a7b2398fac0ac51653ba62_images_image-7f271449-0adb-4af0-83b9-121da8722239.png",
        ASSETS
        / "c__Users_lirik_AppData_Roaming_Cursor_User_workspaceStorage_001952add7a7b2398fac0ac51653ba62_images_image-7c0f8cd0-d149-4f2b-846c-d200c2730bc5.png",
        ASSETS
        / "c__Users_lirik_AppData_Roaming_Cursor_User_workspaceStorage_001952add7a7b2398fac0ac51653ba62_images_image-ca8d3f46-4206-4536-8733-d625b5c7ffde.png",
    ]
    for p in shots:
        if not p.exists():
            raise SystemExit(f"Не найден скриншот: {p}")

    plantuml = ROOT / "db_structure" / "plantuml.png"
    if not plantuml.exists():
        raise SystemExit(f"Не найден файл диаграммы: {plantuml}")

    chart_path = REPORT / "presentation_recommender_quality.png"
    _build_quality_chart(chart_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def S(title: str, body: list[str], notes: str) -> None:
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        _set_slide_dark_bg(slide)
        _add_title(slide, title)
        _add_body(slide, body)
        _speaker(slide, notes)

    def IMG(title: str, img: Path, notes: str, img_height: float = 5.45) -> None:
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        _set_slide_dark_bg(slide)
        _add_title(slide, title, top=0.28, height=0.78, size_pt=IMG_TITLE_PT)
        pic = slide.shapes.add_picture(str(img), Inches(0.45), Inches(0.98), height=Inches(img_height))
        slide_w = float(prs.slide_width)
        pic_w = pic.width
        pic.left = int((slide_w - pic_w) / 2)
        _speaker(slide, notes)

    # 1 — титул
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_dark_bg(slide)
    tb = slide.shapes.add_textbox(Inches(0.65), Inches(2.2), Inches(12.0), Inches(3.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Аналитика футбольных клубов"
    p.font.name = "Segoe UI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TITLE
    p2 = tf.add_paragraph()
    p2.text = "PostgreSQL · Python · Streamlit"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(26)
    p2.font.color.rgb = ACCENT
    p2.space_before = Pt(16)
    _speaker(
        slide,
        "Около 7 минут: зачем система, как устроена, рекомендации в SQL, три экрана, итог.",
    )

    S(
        "Задача и содержание",
        [
            "Каталог клубов и игроков с фильтрами и отчётами.",
            "Скаут: отбор по полям БД, пагинация, сортировка.",
            "Рекомендации: те же фильтры, порядок по расстоянию до заданного профиля (SQL).",
        ],
        "Три пункта без лишних слов: данные, поиск, рекомендации.",
    )

    S(
        "Архитектура",
        [
            "БД: PostgreSQL, скрипты sql_scripts/ddl.sql и dml.sql.",
            "Запросы: db_connection.py (psycopg2, параметры).",
            "Интерфейс: pages.py; вход и роли: auth.py.",
            "Фото игроков: скрипт tools/upload_player_photos_to_s3.py, в таблице хранится URL.",
        ],
        "Схему «клиент — приложение — БД» можно показать на доске за 20 секунд.",
    )

    IMG(
        "Модель данных (ER)",
        plantuml,
        "Не зачитывать все таблицы: Club–Player, Player–Country, подтип игрока, трофеи.",
        img_height=5.55,
    )

    S(
        "Сущности и роли",
        [
            "Клубы, стадионы, города, страны; игроки, статистика, амплуа.",
            "Пользователь: просмотр и фильтры.",
            "Репортёр: добавление трофея клубу.",
            "Администратор: перевод игрока в другой клуб и зарплата.",
        ],
        "Роли — для разных сценариев изменения данных.",
    )

    # Один слайд — три скрина (после enrich вставятся 2 слайда со схемами)
    layout = prs.slide_layouts[6]
    slide_ui = prs.slides.add_slide(layout)
    _set_slide_dark_bg(slide_ui)
    _add_title(slide_ui, "Интерфейс: три экрана", top=0.28, height=0.72, size_pt=IMG_TITLE_PT)
    _combine_three_screenshots(slide_ui, shots)
    _speaker(
        slide_ui,
        "Слева направо: фильтр клубов, список игроков клуба, форма цели и весов для рекомендаций.",
    )

    S(
        "Скаут и рекомендации",
        [
            "Скаут: возраст, клуб, амплуа, голы и передачи за сезон, пагинация.",
            "Рекомендации: вариант B — взвешенная сумма отклонений от цели по возрасту, зарплате, голам, передачам.",
            "Сначала тот же отбор кандидатов, что у скаута; сортировка по возрастанию метрики.",
        ],
        "Подчеркнуть: метрика считается в SQL (search_ideal_profile_players), не вручную в Python.",
    )

    S(
        "Формула расстояния (как в проекте)",
        [
            "d = wв·|возраст−цель|/40 + wз·|зарплата−цель|/50 000 000",
            "  + wг·|голы−цель|/30 + wп·|передачи−цель|/25",
            "Нормировки под масштаб полей; NULL в статистике — как 0 (COALESCE в SQL).",
        ],
        "Коэффициенты задаёт пользователь на форме; формула видна комиссии.",
    )

    IMG(
        "График: порядок кандидатов по d",
        chart_path,
        "Те же числа, что в SQL: меньше d — ближе к целевому профилю.",
        img_height=5.15,
    )

    S(
        "Фото, безопасность",
        [
            "Фото: Sports.ru / Wikipedia → при успехе запись URL; UI не ходит на внешние сайты при каждом показе.",
            "Запросы с подстановкой параметров; пароли — хеш; строка подключения из .env.",
        ],
        "Коротко про эксплуатацию, без маркетинга.",
    )

    S(
        "Что сделали и кому помогли",
        [
            "Сделали рабочую систему: база данных + веб-интерфейс + рекомендации игроков.",
            "Помогли скауту: меньше ручного просмотра, быстрее первичный отбор кандидатов.",
            "Помогли аналитику и тренеру: единая картина по клубам и игрокам на одном экране.",
            "Итог для клуба: решение по кандидатам принимается быстрее и на данных, а не «на глаз».",
        ],
        "Это слайд-ответ на вопрос «Ну и что?»: перечислите результат и прямую пользу для ролей.",
    )

    S(
        "Итог",
        [
            "Нормализованная БД и сценарии в Streamlit.",
            "Рекомендации — явная метрика и сортировка в SQL.",
            "Дальше по желанию: тесты, журнал операций, отдельный слой API.",
        ],
        "Один абзац: что сделано и что логично добавить.",
    )

    out = REPORT / "presentation.pptx"
    prs.save(str(out))
    print(f"Сохранено: {out} ({len(prs.slides)} слайдов). Затем: python tools/enrich_presentation_diagrams.py")


if __name__ == "__main__":
    main()
